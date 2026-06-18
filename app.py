"""
ETV Auto Sticker Dashboard - Production Version
Google Vision API for OCR (zero server RAM)
MongoDB Atlas for storage
Render Starter for hosting
"""
from flask import Flask, request, jsonify, send_from_directory, send_file
import os, re, hashlib, base64, io, json, requests as req_lib
from pathlib import Path
from datetime import datetime
from collections import defaultdict

app = Flask(__name__, static_folder=os.path.join(os.path.dirname(os.path.abspath(__file__)), "static"))

UPLOAD_FOLDER = Path("uploads")
OUTPUT_FOLDER = Path("output")
UPLOAD_FOLDER.mkdir(exist_ok=True)
OUTPUT_FOLDER.mkdir(exist_ok=True)

SUPPORTED_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".bmp"}
GOOGLE_VISION_KEY = os.environ.get("GOOGLE_VISION_KEY", "")
MONGO_URI = os.environ.get("MONGO_URI", "")

# ── MongoDB ──────────────────────────────────────────────────
_db = None
def get_db():
    global _db
    if _db is None and MONGO_URI:
        from pymongo import MongoClient
        _db = MongoClient(MONGO_URI, serverSelectionTimeoutMS=5000)["etv_dashboard"]["photos"]
    return _db

def save_record(r):
    col = get_db()
    if col is not None:
        col.update_one({"id": r["id"]}, {"$set": r}, upsert=True)

def load_records():
    col = get_db()
    return list(col.find({}, {"_id": 0})) if col is not None else []

def record_exists(pid):
    col = get_db()
    return col.find_one({"id": pid}) is not None if col is not None else False

# ── Series map ───────────────────────────────────────────────
SERIES_MAP = {
  "A":{"location":"VISHAKAPATNAM","state":"Andhra Pradesh","qty":250},
  "B":{"location":"ANAKAPALLI","state":"Andhra Pradesh","qty":13},
  "C":{"location":"VIJAYAWADA","state":"Andhra Pradesh","qty":150},
  "D":{"location":"GUDIWADA","state":"Andhra Pradesh","qty":16},
  "E":{"location":"Machilipatnam","state":"Andhra Pradesh","qty":40},
  "F":{"location":"GUNTUR","state":"Andhra Pradesh","qty":150},
  "G":{"location":"NARSARAOPET","state":"Andhra Pradesh","qty":21},
  "H":{"location":"TENALI","state":"Andhra Pradesh","qty":30},
  "I":{"location":"Mangalagiri","state":"Andhra Pradesh","qty":50},
  "J":{"location":"NELLORE","state":"Andhra Pradesh","qty":60},
  "K":{"location":"KAVALI","state":"Andhra Pradesh","qty":16},
  "L":{"location":"ELURU","state":"Andhra Pradesh","qty":45},
  "M":{"location":"BHIMAVARAM","state":"Andhra Pradesh","qty":28},
  "N":{"location":"TANUKU","state":"Andhra Pradesh","qty":14},
  "O":{"location":"TADEPALLIGUDEM","state":"Andhra Pradesh","qty":19},
  "P":{"location":"Dubacharla","state":"Andhra Pradesh","qty":4},
  "Q":{"location":"Athili","state":"Andhra Pradesh","qty":6},
  "R":{"location":"Narasapuram","state":"Andhra Pradesh","qty":20},
  "S":{"location":"Yarnagudem","state":"Andhra Pradesh","qty":7},
  "T":{"location":"Palakollu","state":"Andhra Pradesh","qty":22},
  "U":{"location":"Gokavaram","state":"Andhra Pradesh","qty":6},
  "V":{"location":"Nidadavolu","state":"Andhra Pradesh","qty":14},
  "W":{"location":"Rajahmundry-Bypass","state":"Andhra Pradesh","qty":12},
  "X":{"location":"Rajahmundry-Gokavaram","state":"Andhra Pradesh","qty":14},
  "Y":{"location":"Rajahmundry-LalaCheruvu","state":"Andhra Pradesh","qty":11},
  "Z":{"location":"Rajahmundry-NandamGaniraju","state":"Andhra Pradesh","qty":10},
  "AA":{"location":"Rajahmundry-ShyamalaTheater","state":"Andhra Pradesh","qty":10},
  "AB":{"location":"Rajahmundry-ThaditotaMarket","state":"Andhra Pradesh","qty":13},
  "AC":{"location":"Seetanagaram","state":"Andhra Pradesh","qty":5},
  "AD":{"location":"Thallapudi","state":"Andhra Pradesh","qty":5},
  "AE":{"location":"Undrajavaram","state":"Andhra Pradesh","qty":4},
  "AF":{"location":"Vadapally","state":"Andhra Pradesh","qty":5},
  "AG":{"location":"Kolamuru","state":"Andhra Pradesh","qty":4},
  "AH":{"location":"Kovvuru","state":"Andhra Pradesh","qty":12},
  "AI":{"location":"Korukonda","state":"Andhra Pradesh","qty":6},
  "AJ":{"location":"Chagallu","state":"Andhra Pradesh","qty":6},
  "AK":{"location":"Amalapuram","state":"Andhra Pradesh","qty":18},
  "AL":{"location":"Ambajipeta","state":"Andhra Pradesh","qty":5},
  "AM":{"location":"Kothapeta","state":"Andhra Pradesh","qty":9},
  "AN":{"location":"Pulletikuru","state":"Andhra Pradesh","qty":3},
  "AO":{"location":"Ravulapalem","state":"Andhra Pradesh","qty":14},
  "AP":{"location":"OldGunturRoad","state":"Andhra Pradesh","qty":12},
  "AQ":{"location":"Peddapuram","state":"Andhra Pradesh","qty":14},
  "AR":{"location":"Pithapuram","state":"Andhra Pradesh","qty":14},
  "AS":{"location":"Samarlakota","state":"Andhra Pradesh","qty":16},
  "AT":{"location":"Kakinada","state":"Andhra Pradesh","qty":90},
  "AU":{"location":"Tuni","state":"Andhra Pradesh","qty":22},
  "AV":{"location":"Jaggampeta","state":"Andhra Pradesh","qty":8},
  "AW":{"location":"Bhimadolu","state":"Andhra Pradesh","qty":5},
  "AX":{"location":"DwarakaTirumala","state":"Andhra Pradesh","qty":8},
  "AY":{"location":"Polavaram","state":"Andhra Pradesh","qty":5},
  "AZ":{"location":"Nuzvid","state":"Andhra Pradesh","qty":18},
  "BA":{"location":"ONGOLE","state":"Andhra Pradesh","qty":32},
  "BB":{"location":"CHIRALA","state":"Andhra Pradesh","qty":27},
  "BC":{"location":"Markapur","state":"Andhra Pradesh","qty":18},
  "BD":{"location":"SRIKAKULAM","state":"Andhra Pradesh","qty":24},
  "BE":{"location":"Vizianagaram","state":"Andhra Pradesh","qty":65},
  "BF":{"location":"Bapatla","state":"Andhra Pradesh","qty":20},
  "BG":{"location":"Sattenapalle","state":"Andhra Pradesh","qty":18},
  "BH":{"location":"Raghavapuram","state":"Andhra Pradesh","qty":3},
  "BI":{"location":"Secunderabad","state":"Telangana","qty":25},
  "BJ":{"location":"Bowenpally","state":"Telangana","qty":25},
  "BK":{"location":"Karkhana","state":"Telangana","qty":22},
  "BL":{"location":"Lingampally","state":"Telangana","qty":25},
  "BM":{"location":"Shamsabad","state":"Telangana","qty":25},
  "BN":{"location":"Bandlaguda","state":"Telangana","qty":22},
  "BO":{"location":"Shaapur","state":"Telangana","qty":15},
  "BP":{"location":"BalajiNagar","state":"Telangana","qty":12},
  "BQ":{"location":"Mallapur","state":"Telangana","qty":20},
  "BR":{"location":"Nizampet","state":"Telangana","qty":25},
  "BS":{"location":"Chintal","state":"Telangana","qty":18},
  "BT":{"location":"ECIL","state":"Telangana","qty":30},
  "BU":{"location":"Kompally","state":"Telangana","qty":30},
  "BV":{"location":"Medchal","state":"Telangana","qty":25},
  "BW":{"location":"Malkajgiri","state":"Telangana","qty":35},
  "BX":{"location":"Bollaram","state":"Telangana","qty":15},
  "BY":{"location":"Patancheru-BHEL","state":"Telangana","qty":35},
  "BZ":{"location":"Mallepalli","state":"Telangana","qty":18},
  "CA":{"location":"Ramanthapur","state":"Telangana","qty":22},
  "CB":{"location":"Malakpet","state":"Telangana","qty":30},
  "CC":{"location":"Kukatpalli","state":"Telangana","qty":35},
  "CD":{"location":"Ameerpet","state":"Telangana","qty":50},
  "CE":{"location":"Dilsukhnagar","state":"Telangana","qty":50},
  "CF":{"location":"LBNagar","state":"Telangana","qty":49},
  "CG":{"location":"Mehdipatnam","state":"Telangana","qty":45},
  "CH":{"location":"Charminar","state":"Telangana","qty":50},
  "CI":{"location":"Koti-Abids","state":"Telangana","qty":45},
  "CJ":{"location":"Miyapur","state":"Telangana","qty":40},
  "CK":{"location":"Uppal","state":"Telangana","qty":40},
  "CL":{"location":"Gachibowli","state":"Telangana","qty":35},
  "CM":{"location":"Madhapur","state":"Telangana","qty":35},
  "CN":{"location":"Tarnaka","state":"Telangana","qty":25},
  "CO":{"location":"RahmatNagar","state":"Telangana","qty":20},
  "CP":{"location":"LakdiKaPul","state":"Telangana","qty":30},
  "CQ":{"location":"WARANGAL","state":"Telangana","qty":120},
  "CR":{"location":"Kazipet","state":"Telangana","qty":18},
  "CS":{"location":"KARIMNAGAR","state":"Telangana","qty":50},
  "CT":{"location":"JAGITYAL","state":"Telangana","qty":16},
  "CU":{"location":"KHAMMAM","state":"Telangana","qty":50},
  "CV":{"location":"PALLONCHA","state":"Telangana","qty":14},
  "CW":{"location":"Sathupally","state":"Telangana","qty":18},
  "CX":{"location":"NIZAMABAD","state":"Telangana","qty":57},
  "CY":{"location":"KAMAREDDY","state":"Telangana","qty":14},
  "CZ":{"location":"Armoor","state":"Telangana","qty":18},
  "DA":{"location":"Bodhan","state":"Telangana","qty":18},
  "DB":{"location":"ADILABAD","state":"Telangana","qty":24},
  "DC":{"location":"MANCHERIAL","state":"Telangana","qty":14},
  "DD":{"location":"Jangaon","state":"Telangana","qty":18},
  "DE":{"location":"Bhongir","state":"Telangana","qty":22},
  "DF":{"location":"NALGONDA","state":"Telangana","qty":25},
  "DG":{"location":"SURYAPET","state":"Telangana","qty":19},
  "DH":{"location":"MIRYALAGUDA","state":"Telangana","qty":18},
  "DI":{"location":"MAHABOOBNAGAR","state":"Telangana","qty":35},
  "DJ":{"location":"Tandur","state":"Telangana","qty":22},
  "DK":{"location":"Vikarabad","state":"Telangana","qty":20},
  "DL":{"location":"Gadwal","state":"Telangana","qty":18},
  "DM":{"location":"Wanaparthy","state":"Telangana","qty":18},
  "DN":{"location":"SIDDIPET","state":"Telangana","qty":19},
  "DO":{"location":"Bhadrachalam","state":"Telangana","qty":22},
  "DP":{"location":"Manuguru","state":"Telangana","qty":16},
  "DQ":{"location":"KURNOOL","state":"Andhra Pradesh","qty":75},
  "DR":{"location":"ADONI","state":"Andhra Pradesh","qty":24},
  "DS":{"location":"YEMMIGANUR","state":"Andhra Pradesh","qty":14},
  "DT":{"location":"Tadipatri","state":"Andhra Pradesh","qty":25},
  "DU":{"location":"Hindupur","state":"Andhra Pradesh","qty":35},
  "DV":{"location":"Dharmavaram","state":"Andhra Pradesh","qty":30},
  "DW":{"location":"ANANTAPURAM","state":"Andhra Pradesh","qty":56},
  "DX":{"location":"GUNTAKAL","state":"Andhra Pradesh","qty":19},
  "DY":{"location":"KADIRI","state":"Andhra Pradesh","qty":13},
  "DZ":{"location":"PRODDUTUR","state":"Andhra Pradesh","qty":30},
  "EA":{"location":"RAYACHOTI","state":"Andhra Pradesh","qty":13},
  "EB":{"location":"KADAPA","state":"Andhra Pradesh","qty":54},
  "EC":{"location":"TIRUPATHI","state":"Andhra Pradesh","qty":60},
  "ED":{"location":"CHITOOR","state":"Andhra Pradesh","qty":27},
  "EE":{"location":"MADANAPALLY","state":"Andhra Pradesh","qty":25},
  "EF":{"location":"Nandyal","state":"Andhra Pradesh","qty":55},
}

DESIGNS = ["Aaro_Pranam","Ammoru","Andhal_Rakshasi","Bommarillu","Janaki_Parinayam",
           "Jhansi","Manasantha_Nuvve","Merupu_Kalalu","Rangula_Ratnam","Sandhya_Ragam",
           "Vasundhara","Veyi_Shubhamulu_Kaluguniku","Yashodha"]

# Telugu + English keywords for each show — Google Vision reads Telugu from sticker
SHOW_KEYWORDS = {
    "Aaro_Pranam":                ["ఆరో ప్రాణం", "aaro pranam", "aro pranam", "ప్రాణం"],
    "Ammoru":                     ["అమ్మోరు", "ammoru"],
    "Andhal_Rakshasi":            ["అందాల రాక్షసి", "andhal rakshasi", "andala", "రాక్షసి"],
    "Bommarillu":                 ["బొమ్మరిల్లు", "bommarillu", "బొమ్మరి"],
    "Janaki_Parinayam":           ["జానకి పరిణయం", "janaki parinayam", "parinayam", "పరిణయం"],
    "Jhansi":                     ["ఝాన్సీ", "jhansi", "ఝాన్"],
    "Manasantha_Nuvve":           ["మనసంతా నువ్వే", "manasantha nuvve", "manasanta", "మనసంతా"],
    "Merupu_Kalalu":              ["మెరుపు కలలు", "merupu kalalu", "మెరుపు"],
    "Rangula_Ratnam":             ["రంగుల రత్నం", "rangula ratnam", "రంగుల", "రత్నం"],
    "Sandhya_Ragam":              ["సంధ్యారాగం", "sandhya ragam", "సంధ్యా", "sandya"],
    "Vasundhara":                 ["వసుంధర", "vasundhara", "వసుంధ"],
    "Veyi_Shubhamulu_Kaluguniku": ["వేయి శుభాలు", "veyi shubhamulu", "శుభాలు", "kaluguniku", "కలుగు"],
    "Yashodha":                   ["యశోద", "yashodha", "yashoda", "యశో"],
}

def detect_design_from_ocr(ocr_text):
    """
    Match show name from OCR text — works with Telugu and English.
    Also tries partial Telugu character matching for robustness.
    """
    if not ocr_text:
        return None
    # Direct keyword match
    for design, keywords in SHOW_KEYWORDS.items():
        for kw in keywords:
            if kw.lower() in ocr_text.lower() or kw in ocr_text:
                return design
    # Partial match — check if any 3+ char substring matches
    for design, keywords in SHOW_KEYWORDS.items():
        for kw in keywords:
            if len(kw) >= 4:
                # Check any 4-char window
                for i in range(len(kw)-3):
                    chunk = kw[i:i+4]
                    if chunk in ocr_text:
                        return design
    return None

GPS_KEYWORDS = {
    "guntur":"GUNTUR","vijayawada":"VIJAYAWADA","visakhapatnam":"VISHAKAPATNAM",
    "vizag":"VISHAKAPATNAM","nellore":"NELLORE","kakinada":"Kakinada",
    "rajahmundry":"Rajahmundry","eluru":"ELURU","warangal":"WARANGAL",
    "karimnagar":"KARIMNAGAR","khammam":"KHAMMAM","nizamabad":"NIZAMABAD",
    "hyderabad":"Hyderabad","secunderabad":"Secunderabad","tirupati":"TIRUPATHI",
    "tirupathi":"TIRUPATHI","kurnool":"KURNOOL","kadapa":"KADAPA",
    "anantapur":"ANANTAPURAM","anantapuram":"ANANTAPURAM","srikakulam":"SRIKAKULAM",
    "vizianagaram":"Vizianagaram","ongole":"ONGOLE","nalgonda":"NALGONDA",
    "mahaboobnagar":"MAHABOOBNAGAR","sambasiva":"GUNTUR","narasaraopet":"NARSARAOPET",
    "tenali":"TENALI","bhimavaram":"BHIMAVARAM","tanuku":"TANUKU","chirala":"CHIRALA",
    "nandyal":"Nandyal","kadiri":"KADIRI","guntakal":"GUNTAKAL","adoni":"ADONI",
    "hindupur":"Hindupur","dharmavaram":"Dharmavaram","tadipatri":"Tadipatri",
    "machilipatnam":"Machilipatnam","gudiwada":"GUDIWADA","kavali":"KAVALI",
    "ongole":"ONGOLE","markapur":"Markapur","bapatla":"Bapatla",
    "mangalagiri":"Mangalagiri","tadepalli":"Mangalagiri",
}

# ── Google Vision OCR ────────────────────────────────────────
def vision_ocr(image_bytes):
    """Send image to Google Vision API, return all detected text."""
    if not GOOGLE_VISION_KEY:
        return ""
    try:
        b64 = base64.b64encode(image_bytes).decode()
        payload = {
            "requests": [{
                "image": {"content": b64},
                "features": [{"type": "TEXT_DETECTION", "maxResults": 1}]
            }]
        }
        resp = req_lib.post(
            f"https://vision.googleapis.com/v1/images:annotate?key={GOOGLE_VISION_KEY}",
            json=payload, timeout=10
        )
        data = resp.json()
        annotations = data.get("responses", [{}])[0].get("textAnnotations", [])
        return annotations[0].get("description", "") if annotations else ""
    except Exception as e:
        print(f"Vision API error: {e}")
        return ""

def extract_gps_region(image_bytes):
    """Crop bottom 28% of image — GPS stamp area."""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        w, h = img.size
        crop = img.crop((0, int(h * 0.72), w, h))
        buf = io.BytesIO()
        crop.save(buf, format="JPEG", quality=85)
        return buf.getvalue()
    except:
        return image_bytes

def extract_sticker_region(image_bytes):
    """
    Crop precisely to sticker area where F-number is written.
    F-number is handwritten in a circle on the right side of the sticker,
    between the show title and timings box — approx x:35-75%, y:17-55%.
    This excludes the number plate zone entirely.
    """
    try:
        from PIL import Image, ImageEnhance
        img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        w, h = img.size
        # Sticker bounds: x 7-76%, y 17-58%
        # F-number specifically: right-center of sticker
        # Crop entire sticker but stop well above number plate
        crop = img.crop((int(w*0.07), int(h*0.17), int(w*0.76), int(h*0.57)))
        # Enhance for handwriting detection
        crop = ImageEnhance.Contrast(crop).enhance(2.5)
        crop = ImageEnhance.Sharpness(crop).enhance(2.0)
        buf = io.BytesIO()
        crop.save(buf, format="JPEG", quality=92)
        return buf.getvalue()
    except:
        return image_bytes

def extract_number_plate(ocr_text):
    """
    Extract Indian vehicle number plate from full-image OCR text.
    Handles: AP39 TC2515, GNT TR 6782, TS07 EF1234 etc.
    """
    STATE_CODES = r'(?:AP|TS|TN|KA|MH|DL|GJ|RJ|UP|MP|WB|OR|HR|PB|HP|JK|BR|JH|CG|AS|KL|GNT)'
    patterns = [
        rf'\b({STATE_CODES})\s*(\d{{1,2}})\s*([A-Z]{{1,2}})\s*(\d{{4}})\b',
        rf'\b(GNT)\s+([A-Z]{{1,2}})\s*(\d{{3,4}})\b',
    ]
    for pat in patterns:
        m = re.search(pat, ocr_text, re.IGNORECASE)
        if m:
            groups = [g for g in m.groups() if g]
            return ' '.join(groups).upper()
    return None

def detect_location(ocr_text):
    """Detect city from OCR text using GPS keywords."""
    text_lower = ocr_text.lower()
    for keyword, city in GPS_KEYWORDS.items():
        if keyword in text_lower:
            return city
    return None

def detect_series_from_location(location):
    """Reverse lookup — location name → series code."""
    if not location:
        return None
    loc_upper = location.upper()
    for series, info in SERIES_MAP.items():
        if info["location"].upper() == loc_upper:
            return series
    return None

def extract_f_number(ocr_text):
    """
    Extract F-number from sticker OCR.
    Only accepts series codes from our 136-centre map.
    Strictly filters number plates, large numbers, and non-campaign codes.
    """
    NUMBER_PLATE_CODES = {
        "AP","TS","TN","KA","MH","DL","GJ","RJ","UP","MP",
        "WB","OR","HR","PB","HP","JK","BR","JH","CG","AS",
        "GNT","HYD","KL","GA","TR","NL"
    }
    candidates = []
    for m in re.finditer(r'\b([A-Z]{1,2})\s*(\d{1,3})\b', ocr_text, re.IGNORECASE):
        series = m.group(1).upper()
        number = int(m.group(2))
        # Must be in our series map AND not a number plate code
        if series in SERIES_MAP and series not in NUMBER_PLATE_CODES:
            if 1 <= number <= 500:  # valid F-number range
                candidates.append((series, f"{series}{number}", number))
    if not candidates:
        return None, None
    candidates.sort(key=lambda x: x[2])
    return candidates[0][0], candidates[0][1]

def extract_datetime(ocr_text):
    """Extract date/time from GPS stamp text."""
    patterns = [
        r'(\d{4}[-/]\d{2}[-/]\d{2}[^\n]{0,25}\d{2}:\d{2})',
        r'(\d{2}[-/]\d{2}[-/]\d{4}[^\n]{0,25}\d{2}:\d{2})',
        r'(\d{4}-\d{2}-\d{2})',
    ]
    for pat in patterns:
        m = re.search(pat, ocr_text)
        if m:
            return m.group(1)
    return ""

def compress_for_storage(image_bytes, max_size=1400):
    """
    Store high-res image for zoom quality.
    1400px allows zooming into sticker details clearly.
    ~300KB per photo, well within MongoDB 16MB document limit.
    """
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        img.thumbnail((max_size, max_size))
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=82)
        return buf.getvalue()
    except:
        return image_bytes[:1000000]

def image_hash(data):
    return hashlib.sha256(data).hexdigest()[:12]

# ── Main photo processor ─────────────────────────────────────
def process_photo(file_bytes, filename, reviewer="Team"):
    """
    Full auto-processing using Google Vision.
    No manual input required from team.
    """
    record = {
        "id": image_hash(file_bytes),
        "filename": filename,
        "reviewer": reviewer,
        "uploaded_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
        "series": None, "location": None, "state": None,
        "design": "—", "f_number": None, "auto_plate": None,
        "gps_text": "", "datetime_stamp": "",
        "status": "Manual Check", "issues": [], "notes": "",
        "image_b64": base64.b64encode(compress_for_storage(file_bytes)).decode()
    }

    # ── Step 1: OCR the GPS stamp (bottom 28%) ────────────────
    gps_bytes  = extract_gps_region(file_bytes)
    gps_text   = vision_ocr(gps_bytes)
    record["gps_text"] = gps_text[:300]

    # Extract datetime from GPS text
    record["datetime_stamp"] = extract_datetime(gps_text)

    # Detect location from GPS text
    location = detect_location(gps_text)

    # ── Step 1b: Extract number plate ───────────────────────────
    # Number plate is in the lower-center of the auto body (y 55-80%)
    # We crop that area precisely and run Vision OCR on it
    try:
        from PIL import Image as PILImg, ImageEnhance as PILEnh
        img_tmp = PILImg.open(io.BytesIO(file_bytes)).convert("RGB")
        w_tmp, h_tmp = img_tmp.size
        # Crop number plate region — lower half, exclude GPS stamp
        plate_region = img_tmp.crop((int(w_tmp*0.05), int(h_tmp*0.56),
                                     int(w_tmp*0.95), int(h_tmp*0.76)))
        # Enhance contrast for better plate reading
        plate_region = PILEnh.Contrast(plate_region).enhance(2.0)
        plate_buf = io.BytesIO()
        plate_region.save(plate_buf, format="JPEG", quality=92)
        plate_ocr = vision_ocr(plate_buf.getvalue())
        auto_plate = extract_number_plate(plate_ocr)
        # Fallback: search in GPS text (sometimes plate appears there)
        if not auto_plate:
            auto_plate = extract_number_plate(gps_text)
        # Fallback: search in full sticker OCR
        if not auto_plate and sticker_text:
            auto_plate = extract_number_plate(sticker_text)
        record["auto_plate"] = auto_plate
    except Exception as e:
        record["auto_plate"] = None
        print(f"Plate extraction error: {e}")

    # ── Step 2: OCR the sticker area ────────────────────────────
    sticker_bytes = extract_sticker_region(file_bytes)
    sticker_text  = vision_ocr(sticker_bytes)

    # Detect design from show title text on sticker
    # Also try with GPS text in case sticker OCR missed it
    detected_design = detect_design_from_ocr(sticker_text)
    if not detected_design:
        # Try combining both OCR results
        detected_design = detect_design_from_ocr(sticker_text + " " + gps_text)
    if detected_design:
        record["design"] = detected_design
        print(f"Design detected: {detected_design}")

    # Extract F-number and series from sticker
    series, f_number = extract_f_number(sticker_text)
    record["series"]   = series
    record["f_number"] = f_number

    # If GPS didn't give location, try series lookup from sticker OCR
    if not location and series and series in SERIES_MAP:
        info = SERIES_MAP.get(series, {})
        location = info.get("location")

    record["location"] = location or "Unknown"

    # IMPORTANT: Series must come from location, not from number plate OCR
    # Look up the correct series from the detected location
    correct_series = None
    for s, info in SERIES_MAP.items():
        if info["location"].upper() == (location or "").upper():
            correct_series = s
            break
    if correct_series:
        record["series"] = correct_series
        record["state"]  = SERIES_MAP[correct_series]["state"]
    elif series and series in SERIES_MAP:
        record["series"] = series
        record["state"]  = SERIES_MAP[series]["state"]

    # ── Step 3: Basic QC ──────────────────────────────────────
    if not location:
        record["issues"].append("Location not detected from GPS stamp")
    if not gps_text.strip():
        record["issues"].append("GPS stamp not found in photo")
    if not detected_design:
        record["issues"].append("Design not identified — please select manually")
    # Note: F-number is entered manually by team, not blocking auto-approve

    # ── Step 4: Auto-approve if clean ────────────────────────
    if not record["issues"]:
        record["status"] = "Approved"
    else:
        record["status"] = "Manual Check"

    return record

# ── Routes ────────────────────────────────────────────────────
@app.route("/")
def index():
    static_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "static")
    return send_from_directory(static_dir, "index.html")

@app.route("/series_map")
def series_map_route():
    return jsonify(SERIES_MAP)

@app.route("/designs")
def designs_route():
    return jsonify(DESIGNS)

@app.route("/upload", methods=["POST"])
def upload():
    files    = request.files.getlist("photos")
    reviewer = request.form.get("reviewer", "Team")
    # Optional manual overrides (used only if Vision fails)
    manual_series  = request.form.get("series", "")
    manual_design  = request.form.get("design", "")
    results = []

    for f in files:
        if not f.filename: continue
        if Path(f.filename).suffix.lower() not in SUPPORTED_EXTS: continue

        file_bytes = f.read()
        pid = image_hash(file_bytes)

        if record_exists(pid):
            results.append({"filename": f.filename, "status": "Duplicate",
                            "id": pid, "issues": ["Already uploaded"]})
            continue

        rec = process_photo(file_bytes, f.filename, reviewer)

        # Apply manual overrides only if auto-detection failed
        if manual_series and not rec["series"]:
            rec["series"] = manual_series.upper()
            info = SERIES_MAP.get(manual_series.upper(), {})
            if info:
                rec["location"] = info["location"]
                rec["state"]    = info["state"]
                rec["issues"]   = [i for i in rec["issues"] if "Location" not in i]

        if manual_design:
            rec["design"] = manual_design

        # Re-evaluate status after overrides
        if rec["issues"]:
            rec["status"] = "Manual Check"
        else:
            rec["status"] = "Approved"

        save_record(rec)
        results.append({k: v for k, v in rec.items() if k != "image_b64"})

    return jsonify({"processed": len(results), "results": results})

@app.route("/records")
def get_records():
    return jsonify([{k: v for k, v in r.items() if k != "image_b64"} for r in load_records()])

@app.route("/stats")
def get_stats():
    records = load_records()
    by_loc  = defaultdict(lambda: {"approved": 0, "manual": 0, "rejected": 0, "total": 0})
    for r in records:
        loc = r.get("location", "Unknown")
        by_loc[loc]["total"] += 1
        s = r.get("status", "")
        if s == "Approved":      by_loc[loc]["approved"] += 1
        elif s == "Manual Check": by_loc[loc]["manual"]  += 1
        elif s == "Rejected":    by_loc[loc]["rejected"] += 1
    return jsonify({
        "total":     len(records),
        "approved":  sum(1 for r in records if r.get("status") == "Approved"),
        "manual":    sum(1 for r in records if r.get("status") == "Manual Check"),
        "rejected":  sum(1 for r in records if r.get("status") == "Rejected"),
        "duplicate": sum(1 for r in records if r.get("status") == "Duplicate"),
        "by_location": dict(by_loc)
    })

@app.route("/location_progress")
def location_progress():
    records = load_records()
    progress = {}
    for series, info in SERIES_MAP.items():
        planned  = info["qty"] * 13
        received = sum(1 for r in records if r.get("series") == series and r.get("status") != "Duplicate")
        approved = sum(1 for r in records if r.get("series") == series and r.get("status") == "Approved")
        progress[series] = {
            "location": info["location"],
            "state":    info["state"],
            "planned":  planned,
            "received": received,
            "approved": approved,
            "percent":  round(received / planned * 100) if planned else 0,
        }
    return jsonify(progress)

@app.route("/image/<photo_id>")
def serve_image(photo_id):
    try:
        col = get_db()
        if col is not None:
            rec = col.find_one({"id": photo_id}, {"image_b64": 1})
            if rec and rec.get("image_b64"):
                img_bytes = base64.b64decode(rec["image_b64"])
                buf = io.BytesIO(img_bytes)
                buf.seek(0)
                from flask import Response
                return Response(buf.getvalue(), mimetype="image/jpeg",
                               headers={"Cache-Control": "public, max-age=3600",
                                        "Access-Control-Allow-Origin": "*"})
    except Exception as e:
        print(f"serve_image error: {e}")
    return "Not found", 404

@app.route("/update", methods=["POST"])
def update():
    d   = request.json
    col = get_db()
    if col is not None:
        fields = {k: d[k] for k in ["status", "design", "f_number", "location", "series", "auto_plate", "notes"] if d.get(k) is not None and k in d}
        if fields:
            col.update_one({"id": d["id"]}, {"$set": fields})
    return jsonify({"ok": True})

@app.route("/export_excel")
def export_excel():
    try:
        import pandas as pd
        records = load_records()
        # Friendly column order for Excel
        cols_order = ["uploaded_at","reviewer","location","state","series","f_number",
                      "auto_plate","design","datetime_stamp","status","notes","gps_text","filename","id"]
        df_raw = [{k: v for k, v in r.items() if k != "image_b64"} for r in records]
        df = pd.DataFrame(df_raw)
        # Reorder columns that exist
        existing = [c for c in cols_order if c in df.columns]
        rest = [c for c in df.columns if c not in cols_order]
        df = df[existing + rest]
        out = OUTPUT_FOLDER / "ETV_QC_Report.xlsx"
        df.to_excel(str(out), index=False)
        return send_file(str(out.resolve()), as_attachment=True, download_name="ETV_QC_Report.xlsx")
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/export_ppt")
def export_ppt():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        records = load_records()
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        DARK  = RGBColor(0x1A, 0x1A, 0x2E)
        ACCENT= RGBColor(0xE9, 0x4F, 0x37)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        GOLD  = RGBColor(0xF5, 0xA6, 0x23)
        GREEN = RGBColor(0x27, 0xAE, 0x60)

        def t(sl, tx, l, tp, w, h, sz, bold=False, col=WHITE, al=PP_ALIGN.LEFT):
            tb = sl.shapes.add_textbox(l, tp, w, h)
            tf = tb.text_frame; tf.word_wrap = True
            p  = tf.paragraphs[0]; p.alignment = al
            run = p.add_run(); run.text = str(tx)
            run.font.size = Pt(sz); run.font.bold = bold
            run.font.color.rgb = col; run.font.name = "Calibri"

        s1 = prs.slides.add_slide(blank)
        bg = s1.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
        t(s1, "ETV Auto Sticker Campaign Report", Inches(0.7), Inches(2.0), Inches(12), Inches(1), 38, True, ACCENT, PP_ALIGN.CENTER)
        t(s1, f"{datetime.now().strftime('%d %B %Y')}  |  {len(records)} photos processed",
          Inches(0.7), Inches(3.2), Inches(12), Inches(0.5), 16, col=WHITE, al=PP_ALIGN.CENTER)

        for r in [x for x in records if x.get("status") == "Approved"]:
            sl = prs.slides.add_slide(blank)
            if r.get("image_b64"):
                try:
                    sl.shapes.add_picture(io.BytesIO(base64.b64decode(r["image_b64"])),
                                         Inches(0.3), Inches(0.3), Inches(6.5), Inches(6.9))
                except: pass
            card = sl.shapes.add_shape(1, Inches(7.1), Inches(0.3), Inches(5.9), Inches(6.9))
            card.fill.solid(); card.fill.fore_color.rgb = DARK; card.line.fill.background()
            t(sl, "ETV AUTO STICKER", Inches(7.3), Inches(0.5), Inches(5.5), Inches(0.5), 14, True, ACCENT)
            y = 1.1
            for lbl, val in [
                ("Location",      r.get("location", "—")),
                ("Series / No.",  f"{r.get('series','')} — {r.get('f_number','')}"),
                ("Design",        r.get("design", "—")),
                ("Date / Time",   r.get("datetime_stamp", "—")),
                ("GPS",           (r.get("gps_text") or "—")[:50]),
                ("Reviewer",      r.get("reviewer", "—")),
            ]:
                t(sl, lbl, Inches(7.3), Inches(y),      Inches(5.5), Inches(0.25), 8,  True,  GOLD)
                t(sl, str(val), Inches(7.3), Inches(y+0.25), Inches(5.5), Inches(0.4),  13, False, WHITE)
                y += 0.85
            b = sl.shapes.add_shape(1, Inches(7.3), Inches(6.55), Inches(2.8), Inches(0.45))
            b.fill.solid(); b.fill.fore_color.rgb = GREEN; b.line.fill.background()
            t(sl, "APPROVED", Inches(7.3), Inches(6.57), Inches(2.8), Inches(0.4), 12, True, WHITE, PP_ALIGN.CENTER)

        buf = io.BytesIO(); prs.save(buf); buf.seek(0)
        return send_file(buf, as_attachment=True, download_name="ETV_Campaign_Report.pptx",
                        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    print(f"\n  ETV Dashboard → http://localhost:{port}\n")
    app.run(debug=False, host="0.0.0.0", port=port)
